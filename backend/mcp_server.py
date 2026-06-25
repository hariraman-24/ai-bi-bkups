import sys
import os

ROOT_DIR = os.path.abspath(os.path.join(os.path.dirname(__file__), ".."))
sys.path.append(ROOT_DIR)

from flask import Flask, request, jsonify
from flask_cors import CORS

from db_connection import connect_db
from postgres_connection import connect_postgres
from sql_generator import SQLGenerator
from visualization.charts import generate_chart
from forecasting.forecast import forecast_sales
from intent_engine import detect_intent
from insight_engine import InsightEngine
from dashboard_engine import DashboardEngine
from llm_engine import LLMEngine

import pandas as pd
import pdfplumber
from docx import Document
import re
from decimal import Decimal
import sqlite3
from pptx import Presentation

# ================= PATH =================
STATIC_DIR = os.path.join(ROOT_DIR, "static")
FILES_DIR = os.path.join(ROOT_DIR, "files")

os.makedirs(STATIC_DIR, exist_ok=True)
os.makedirs(FILES_DIR, exist_ok=True)

# ================= APP =================
app = Flask(__name__)
CORS(app)

sql_generator = SQLGenerator()
insight_engine = InsightEngine()
dashboard_engine = DashboardEngine()
llm_engine_service = LLMEngine()

active_file_path = None
document_text = ""
document_pages = []

# ================= NLP =================
def normalize_question(q):

    q = q.lower().strip()

    replacements = {
        "best": "top",
        "highest": "top",
        "most": "top",
        "worst": "lowest",
        "least": "lowest",
        "increase": "trend",
        "decrease": "trend",
        "growth": "trend"
    }

    for k, v in replacements.items():
        q = q.replace(k, v)

    return q
# ================= DOMAIN =================
def is_business_query(q, enable_ai=False):

    q = q.lower().strip()

    # =================================================
    # ✅ BUSINESS KEYWORDS
    # =================================================
    business_keywords = [

        # SALES / FINANCE
        "sale",
        "sales",
        "revenue",
        "profit",
        "income",
        "expense",
        "expenses",
        "amount",
        "price",
        "cost",
        "finance",

        # BUSINESS ENTITIES
        "customer",
        "customers",
        "employee",
        "employees",
        "product",
        "products",
        "region",
        "category",
        "department",

        # ANALYTICS
        "forecast",
        "trend",
        "growth",
        "analysis",
        "analytics",
        "dashboard",
        "summary",
        "report",
        "kpi",
        "performance",

        # COMPARISON / FILTERS
        "top",
        "highest",
        "lowest",
        "best",
        "worst",
        "compare",
        "distribution",
        "share",
        "contribution",
        "average",
        "total",
        "count",
        "maximum",
        "minimum",

        # TIME
        "month",
        "year",
        "quarter",
        "monthly",
        "yearly"
    ]

    # =================================================
    # ❌ NON BUSINESS FILTER
    # =================================================
    non_business_keywords = [
        "movie",
        "weather",
        "cricket",
        "football",
        "politics",
        "news",
        "cinema",
        "actor",
        "song"
    ]

    # ---------------------------------------------
    # BLOCK NON BUSINESS
    # ---------------------------------------------
    if any(
        word in q
        for word in non_business_keywords
    ):
        return False

    # ---------------------------------------------
    # FAST ALLOW BUSINESS
    # ---------------------------------------------
    if any(
        word in q
        for word in business_keywords
    ):
        return True

    # ---------------------------------------------
    # DYNAMIC LLM CHECK
    # ---------------------------------------------
    if enable_ai:
        return llm_engine_service.is_business_query_llm(q)
    return False

# ================= FORMAT =================
def format_text(data):

    formatted = []

    for row in data:

        line = ", ".join(
            f"{k.replace('_',' ').title()}: {v}"
            for k, v in row.items()
        )

        formatted.append(line)

    return formatted

# ================= FILTER =================
def apply_ranking_filter(question, result):

    if not result:
        return result

    numeric_keys = [
        k for k in result[0]
        if isinstance(result[0][k], (int, float))
    ]

    if not numeric_keys:
        return result

    key = numeric_keys[0]

    if "ascending" in question:
        return sorted(result, key=lambda x: x[key])

    if "descending" in question:
        return sorted(result, key=lambda x: x[key], reverse=True)

    match = re.search(r"top\s+(\d+)", question)

    if match:

        n = int(match.group(1))

        return sorted(
            result,
            key=lambda x: x[key],
            reverse=True
        )[:n]

    if "top" in question:
        return [max(result, key=lambda x: x[key])]

    if "lowest" in question:
        return [min(result, key=lambda x: x[key])]

    return result

# ================= DOCUMENT AI =================
def process_document_question(question):

    global document_text
    global document_pages

    q = question.lower().strip()

    if not document_text:

        return {
            "type": "text",
            "data": ["No document loaded"]
        }

    # ================= PAGE DETECTION =================
    page_match = re.search(r"page\s+(\d+)", q)

    current_page_text = ""

    if page_match:

        page_num = int(page_match.group(1)) - 1

        if 0 <= page_num < len(document_pages):

            current_page_text = document_pages[page_num]

        else:

            return {
                "type": "text",
                "data": ["Page not found"]
            }

    # ================= WORD COUNT =================
    if any(x in q for x in [
        "how many words",
        "word count",
        "total words",
        "number of words"
    ]):

        target_text = (
            current_page_text
            if current_page_text
            else document_text
        )

        clean = re.sub(r"\s+", " ", target_text)

        words = re.findall(
            r"\b[a-zA-Z0-9]+\b",
            clean
        )

        if current_page_text:

            return {
                "type": "text",
                "data": [
                    f"Total words in page {page_num + 1}: {len(words)}"
                ]
            }

        return {
            "type": "text",
            "data": [
                f"Total words in document: {len(words)}"
            ]
        }

    # ================= DYNAMIC AI DOCUMENT ROUTING =================
    # To run as fast as possible on local hardware, we carefully control the context length
    # based on the exact intent of the user's question, and then feed it to the AI.
    
    q_lower = q.lower()
    
    # 1. OVERVIEW PATH
    if "overview" in q_lower:
        context = document_text[:800] 
        answer = llm_engine_service.answer_document_question("Provide a concise high-level overview of this document.", context)
        
    # 2. SUMMARIZE PATH
    elif "summary" in q_lower or "summarize" in q_lower:
        # Grab from start, middle, and end for a holistic summary
        total_len = len(document_text)
        if total_len > 1500:
            start = document_text[:500]
            mid = document_text[total_len//2 : total_len//2 + 500]
            end = document_text[-500:]
            context = f"{start}\n...\n{mid}\n...\n{end}"
        else:
            context = document_text
        answer = llm_engine_service.answer_document_question("Summarize the main points of this entire document.", context)
        
    # 3. EXPLANATION / SPECIFIC QUESTION PATH
    else:
        query_words = set(q_lower.split())
        scored_paragraphs = []
        
        for p in document_pages:
            p_lower = p.lower()
            score = sum(1 for word in query_words if word in p_lower and len(word) > 3)
            if score > 0:
                scored_paragraphs.append((score, p))
                
        if scored_paragraphs:
            scored_paragraphs.sort(reverse=True, key=lambda x: x[0])
            best_matches = [p[1] for p in scored_paragraphs[:2]]
            context = "\n\n".join(best_matches)[:1000]
        else:
            context = document_text[:500]
            
        answer = llm_engine_service.answer_document_question(q, context)

    return {
        "type": "text",
        "data": [answer]
    }

# ================= MAIN =================
@app.route("/query", methods=["POST"])
def run_query():

    global active_file_path
    global document_text
    global document_pages

    try:

        # ================= INPUT =================
        if request.files:

            question = normalize_question(
                request.form.get("question", "")
            )

            db_choice = request.form.get(
                "database",
                "mysql"
            ).lower()
            
            enable_ai_str = str(request.form.get("enable_ai", "False"))

        else:

            data = request.json

            question = normalize_question(
                data.get("question", "")
            )

            db_choice = data.get(
                "database",
                "mysql"
            ).lower()
            
            enable_ai_str = str(data.get("enable_ai", "False"))
            
        enable_ai = enable_ai_str.lower() == "true"

        if not question:

            return jsonify({
                "type": "text",
                "data": ["Question required"]
            })

        # ================= INTENT =================
        intent = detect_intent(question)

        if intent == "ranking_top" and "top" not in question:
            question += " top"

        if intent == "ranking_low" and "lowest" not in question:
            question += " lowest"

        # ================= FILE =================
        file = request.files.get("file")

        if file:

            path = os.path.join(
                FILES_DIR,
                file.filename
            )

            file.save(path)

            active_file_path = path

            document_text = ""
            document_pages = []

            filename = active_file_path.lower()

            # ================= CSV =================
            if filename.endswith(".csv"):

                df = pd.read_csv(active_file_path)

                def do_spreadsheet_fallback():
                    df_preview = df.head(15).to_string()
                    ans = llm_engine_service.answer_document_question(question, f"Spreadsheet Data Preview:\n{df_preview}")
                    return jsonify({"type": "chat", "message": ans})

                q_lower = question.lower()
                intercept_words = ["explain", "summary", "summarize", "overview", "about", "what is this", "sheet"]
                if any(w in q_lower for w in intercept_words):
                    return do_spreadsheet_fallback()

                conn_mem = sqlite3.connect(":memory:")
                df.to_sql("uploaded_data", conn_mem, index=False)
                
                schema_context = "uploaded_data(" + ", ".join(df.columns) + ")"
                sql_query = sql_generator.generate_sql(question, schema_context=schema_context, db_type="sqlite")
                
                if not sql_query:
                    return do_spreadsheet_fallback()
                    
                try:
                    result_df = pd.read_sql_query(sql_query, conn_mem)
                    result = result_df.to_dict(orient="records")
                    if not result:
                        return do_spreadsheet_fallback()
                except Exception as e:
                    return do_spreadsheet_fallback()

                result = apply_ranking_filter(
                    question,
                    result
                )

                insight = insight_engine.generate_insight(result)
                if enable_ai:
                    chat_insight = llm_engine_service.generate_chat_insight(question, str(result[:5]))
                    if chat_insight:
                        insight += f"\n\n💡 AI Analysis:\n{chat_insight}"

                dashboard = dashboard_engine.build_dashboard(result)
                
                chart_result = result[:20]
                chart_path, _ = generate_chart(
                    chart_result,
                    question
                    )
                return jsonify({
                    "type": "chart",
                    "chart_path": chart_path,
                    "data": format_text(result[:200]),
                    "summary": dashboard["kpis"],
                    "dashboard": dashboard,
                    "insight": insight
                })

            # ================= XLSX =================
            elif filename.endswith(".xlsx"):

                df = pd.read_excel(active_file_path)

                def do_spreadsheet_fallback():
                    df_preview = df.head(15).to_string()
                    ans = llm_engine_service.answer_document_question(question, f"Spreadsheet Data Preview:\n{df_preview}")
                    return jsonify({"type": "chat", "message": ans})

                q_lower = question.lower()
                intercept_words = ["explain", "summary", "summarize", "overview", "about", "what is this", "sheet"]
                if any(w in q_lower for w in intercept_words):
                    return do_spreadsheet_fallback()

                conn_mem = sqlite3.connect(":memory:")
                df.to_sql("uploaded_data", conn_mem, index=False)
                
                schema_context = "uploaded_data(" + ", ".join(df.columns) + ")"
                sql_query = sql_generator.generate_sql(question, schema_context=schema_context, db_type="sqlite")
                
                if not sql_query:
                    return do_spreadsheet_fallback()
                    
                try:
                    result_df = pd.read_sql_query(sql_query, conn_mem)
                    result = result_df.to_dict(orient="records")
                    if not result:
                        return do_spreadsheet_fallback()
                except Exception as e:
                    return do_spreadsheet_fallback()

                result = apply_ranking_filter(
                    question,
                    result
                )

                insight = insight_engine.generate_insight(result)
                if enable_ai:
                    chat_insight = llm_engine_service.generate_chat_insight(question, str(result[:5]))
                    if chat_insight:
                        insight += f"\n\n💡 AI Analysis:\n{chat_insight}"

                dashboard = dashboard_engine.build_dashboard(result)
                chart_result = result[:20]
                chart_path, _ = generate_chart(
                    chart_result,
                    question
                    )

                
                return jsonify({
                    "type": "chart",
                    "chart_path": chart_path,
                    "data": format_text(result[:200]),
                    "summary": dashboard["kpis"],
                    "dashboard": dashboard,
                    "insight": insight
                })

            # ================= PDF =================
            elif filename.endswith(".pdf"):

                full_text = ""
                pages = []

                with pdfplumber.open(active_file_path) as pdf:

                    for page in pdf.pages:

                        text = page.extract_text() or ""

                        cleaned = text.strip()

                        if cleaned:
                            pages.append(cleaned)

                        full_text += cleaned + "\n"

                document_text = full_text
                document_pages = pages

                response = process_document_question(question)

                return jsonify(response)

            # ================= DOCX =================
            elif filename.endswith((".docx", ".doc")):

                doc = Document(active_file_path)

                paragraphs = [
                    p.text for p in doc.paragraphs
                    if p.text.strip()
                ]

                full_text = "\n".join(paragraphs)

                document_text = full_text

                document_pages = [
                    p.strip()
                    for p in paragraphs
                    if p.strip()
                ]

                response = process_document_question(question)

                return jsonify(response)

            # ================= PPTX / PPT =================
            elif filename.endswith((".pptx", ".ppt")):
                
                try:
                    prs = Presentation(active_file_path)
                    slides_text = []
                    
                    for slide in prs.slides:
                        slide_content = []
                        for shape in slide.shapes:
                            if hasattr(shape, "text"):
                                slide_content.append(shape.text)
                        
                        slides_text.append("\n".join(slide_content))
                        
                    full_text = "\n\n--- Slide ---\n\n".join(slides_text)
                    
                    document_text = full_text
                    document_pages = slides_text
                    
                    response = process_document_question(question)
                    return jsonify(response)
                except Exception as e:
                    return jsonify({
                        "type": "text",
                        "data": [f"Could not read PowerPoint file: {str(e)}"]
                    })
            
            # ================= TXT =================
            elif filename.endswith(".txt"):
                
                with open(active_file_path, "r", encoding="utf-8") as f:
                    full_text = f.read()
                    
                document_text = full_text
                # Splitting into chunks of roughly 500 characters for pages
                chunk_size = 500
                document_pages = [full_text[i:i+chunk_size] for i in range(0, len(full_text), chunk_size)]
                
                response = process_document_question(question)
                return jsonify(response)

            else:

                return jsonify({
                    "type": "text",
                    "data": ["Unsupported file format"]
                })

        # ================= AUTO CLEAR =================
        business_words = [
            "sales", "product", "employee",
            "customer", "finance", "revenue",
            "profit", "forecast", "trend"
        ]

        if any(w in question.lower() for w in business_words):

            document_text = ""
            document_pages = []

        # ================= DOCUMENT =================
        if document_text:

            response = process_document_question(question)

            return jsonify(response)

        # ================= BUSINESS CHECK =================
        if not is_business_query(question, enable_ai):

            answer = llm_engine_service.general_chat(question)
            return jsonify({
                "type": "chat",
                "message": answer
            })

        # ================= FORECAST =================
        if "forecast" in question:

            conn = connect_db()

            if conn is None:

                return jsonify({
                    "type": "text",
                    "data": ["Database connection failed"]
                })

            cursor = conn.cursor()

            cursor.execute("""
                SELECT
                    EXTRACT(YEAR FROM sale_date),
                    EXTRACT(MONTH FROM sale_date),
                    SUM(sales_amount)
                FROM sales
                GROUP BY 1,2
                ORDER BY 1,2
            """)

            rows = cursor.fetchall()

            historical = [
                {
                    "year": int(r[0]),
                    "month": int(r[1]),
                    "total_sales": float(r[2])
                }
                for r in rows
            ]

            result = forecast_sales(historical)

            chart_path, _ = generate_chart(
                result["historical"],
                question
            )

            return jsonify({
                "type": "forecast",
                "chart_path": chart_path,
                "historical": format_text(
                    result["historical"]
                ),
                "forecast": format_text(
                    result["forecast"]
                ),
                "insight": result["insight"]
            })

        # ================= SQL =================
        sql_query = sql_generator.generate_sql(question, db_type=db_choice)

        if not sql_query:

            answer = llm_engine_service.general_chat(question)
            return jsonify({
                "type": "chat",
                "message": answer
            })

        conn = (
            connect_postgres()
            if db_choice == "postgresql"
            else connect_db()
        )

        if conn is None:

            return jsonify({
                "type": "text",
                "data": ["Database connection failed"]
            })

        cursor = conn.cursor()

        cursor.execute(sql_query)

        rows = cursor.fetchall()

        cols = [
            d[0]
            for d in cursor.description
        ]

        result = []

        for row in rows:

            temp = {}

            for c, v in zip(cols, row):

                if isinstance(v, Decimal):
                    temp[c] = float(v)
                else:
                    temp[c] = v

            result.append(temp)

        if not result:

            answer = llm_engine_service.general_chat(question)
            return jsonify({
                "type": "chat",
                "message": answer
            })

        result = apply_ranking_filter(
            question,
            result
        )

        insight = insight_engine.generate_insight(result)
        if enable_ai:
            chat_insight = llm_engine_service.generate_chat_insight(question, str(result[:5]))
            if chat_insight:
                insight += f"\n\n💡 AI Analysis:\n{chat_insight}"

        dashboard = dashboard_engine.build_dashboard(result)
        chart_result = result[:20]
        chart_path, _ = generate_chart(
            chart_result,
            question
            )

        return jsonify({
            "type": "chart",
            "chart_path": chart_path,
            "data": format_text(result[:200]),
            "sql": sql_query,
            "summary": dashboard["kpis"],
            "dashboard": dashboard,
            "insight": insight
        })

    except Exception as e:

        return jsonify({
            "type": "text",
            "data": [str(e)]
        })

# ================= HEALTH =================
@app.route("/health")
def health():

    return jsonify({
        "status": "ok"
    })

# ================= FILE RESET =================
@app.route("/file/detach", methods=["POST"])
def detach_file():

    global active_file_path
    global document_text
    global document_pages

    active_file_path = None
    document_text = ""
    document_pages = []

    return jsonify({
        "message": "File context cleared"
    })

# ================= RUN =================
if __name__ == "__main__":

    app.run(
        debug=True,
        port=5000
    )